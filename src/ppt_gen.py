import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 1. 프레젠테이션 초기화 및 16:9 와이드스크린 설정
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 테마 색상 정의
COLOR_PRIMARY = RGBColor(24, 43, 73)    # 딥네이비 (신뢰감)
COLOR_SECONDARY = RGBColor(70, 130, 180) # 스틸블루 (포인트)
COLOR_TEXT = RGBColor(51, 51, 51)       # 짙은 그레이 (본문)
COLOR_BG_LIGHT = RGBColor(245, 247, 250) # 연한 배경색

def set_font(run, name="맑은 고딕", size=14, bold=False, color=COLOR_TEXT):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color

def add_header(slide, title_text, subtitle_text=""):
    """모든 슬라이드에 적용할 깔끔한 헤더 양식"""
    # 상단 띠 그리기
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_PRIMARY
    shape.line.color.rgb = COLOR_PRIMARY
    
    # 타이틀 텍스트
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title_text
    set_font(run, size=24, bold=True, color=RGBColor(255, 255, 255))
    
    if subtitle_text:
        p2 = tf.add_paragraph()
        run2 = p2.add_run()
        run2.text = subtitle_text
        set_font(run2, size=12, color=RGBColor(200, 220, 240))

# -------------------------------------------------------------
# Slide 1: 타이틀 슬라이드
# -------------------------------------------------------------
blank_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_layout)

# 배경 색상 박스
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
bg.fill.solid()
bg.fill.fore_color.rgb = COLOR_PRIMARY
bg.line.color.rgb = COLOR_PRIMARY

# 대제목
txBox = slide.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(11.333), Inches(2.0))
tf = txBox.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Classic 서비스 Fade-out 및 VPC 이관 전략\n"
set_font(run, size=40, bold=True, color=RGBColor(255, 255, 255))

run_sub = p.add_run()
run_sub.text = "공공 리전 선행 추진 및 글로벌 리전 이슈 대응 방안"
set_font(run_sub, size=20, color=COLOR_SECONDARY)

# 하단 보고 정보
tx_info = slide.shapes.add_textbox(Inches(1.0), Inches(5.8), Inches(5.0), Inches(1.0))
tf_info = tx_info.text_frame
p_info = tf_info.paragraphs[0]
run_info = p_info.add_run()
run_info.text = "보고일자: 2026. 05\n보고대상: 임원진 (수환님 발제 건)"
set_font(run_info, size=12, color=RGBColor(200, 200, 200))


# -------------------------------------------------------------
# Slide 2: 추진 배경 및 전략적 방향성
# -------------------------------------------------------------
slide = prs.slides.add_slide(blank_layout)
add_header(slide, "1. 추진 배경 및 핵심 방향성", "전체 페이드아웃 정체 해소를 위한 KPI 전환 가시화")

# 좌측 카드 (현황 및 문제점)
card1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
card1.fill.solid()
card1.fill.fore_color.rgb = COLOR_BG_LIGHT
card1.line.color.rgb = COLOR_SECONDARY
tf1 = card1.text_frame
tf1.word_wrap = True
tf1.margin_left = Inches(0.3)
tf1.margin_top = Inches(0.3)

p = tf1.paragraphs[0]
r = p.add_run()
r.text = "현재 직면한 교착 상태\n\n"
set_font(r, size=18, bold=True, color=COLOR_PRIMARY)
bullets1 = [
    "사전 준비(신규 VM 차단 등)는 완료되었으나, 사업부 협의 단계에서 중단됨",
    "해외 리전(미국, 독일)의 VPC 미구축으로 인해 글로벌 이관 전면 스탑",
    "대형 고객사(SKT, 벨로프 등)의 특수 네트워크 및 인프라 종속성 존재"
]
for b in bullets1:
    p = tf1.add_paragraph()
    r = p.add_run()
    r.text = "• " + b + "\n"
    set_font(r, size=14)

# 우측 카드 (전략적 제안)
card2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.8))
card2.fill.solid()
card2.fill.fore_color.rgb = COLOR_BG_LIGHT
card2.line.color.rgb = COLOR_PRIMARY
tf2 = card2.text_frame
tf2.word_wrap = True
tf2.margin_left = Inches(0.3)
tf2.margin_top = Inches(0.3)

p = tf2.paragraphs[0]
r = p.add_run()
r.text = "핵심 선회 전략 (Pivot)\n\n"
set_font(r, size=18, bold=True, color=COLOR_PRIMARY)
bullets2 = [
    "KPI 변경 검토: [민간 중심] → [공공 우선 추진]으로 속도 조절",
    "공공 리전은 해외 이슈가 없어 상대적으로 구조가 단순함",
    "공공 선행 전환을 통해 Lesson Learned를 확보한 후 민간 후속 적용",
    "VPC 완전 종료 및 전량 이관을 최종 목적지로 설정"
]
for b in bullets2:
    p = tf2.add_paragraph()
    r = p.add_run()
    r.text = "• " + b + "\n"
    set_font(r, size=14)


# -------------------------------------------------------------
# Slide 3: 기술적 이관 이슈 및 실무적 Bottleneck
# -------------------------------------------------------------
slide = prs.slides.add_slide(blank_layout)
add_header(slide, "2. 5대 기술적 이슈 및 제약사항", "고객 반발 및 기술 차단 요소를 선제적으로 분류")

# 표 생성 (행 6, 열 3)
rows, cols = 6, 3
left, top, width, height = Inches(0.8), Inches(1.8), Inches(11.733), Inches(4.8)
table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
table = table_shape.table

# 컬럼 너비 조절
table.columns[0].width = Inches(2.0)
table.columns[1].width = Inches(4.8)
table.columns[2].width = Inches(4.933)

headers = ["분류", "주요 기술 쟁점 및 현황", "실무적 대응 방안 / 영향도"]
for i, h in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = h
    cell.fill.solid()
    cell.fill.fore_color.rgb = COLOR_PRIMARY
    set_font(cell.text_frame.paragraphs[0], size=14, bold=True, color=RGBColor(255,255,255))

issues = [
    ("OS 레벨", "Classic 고객 대부분이 단종(EOS)된 구버전 OS 사용 중", "VPC 제공 OS 목록과 대조하여 강제 업그레이드 여부 검토 필요"),
    ("네트워크", "Classic → VPC 이전 시 인프라 특성상 IP 전체 변경 불가피", "온프레미스/타 클라우드 연계 고객은 내부 전산 작업 수반 필요"),
    ("SKT 특수성", "SKT 전용 특수 네트워크가 구성되어 있어 VPC 이관 불가", "별도 전용선 연계 방안 수립 또는 영업 부서 협의 통한 예외 처리"),
    ("상품 스펙/비용", "Classic 대비 VPC의 서버 타입, 스펙, 단가 체계 상이", "단가 상승 고객 반발 예상 (ex. 월 1천만 원 → 1.8천만 원 상향 건)"),
    ("보안 서비스", "기존 Classic 보안 상품과 VPC 보안 아키텍처 상이", "보안 조직 협조를 통해 CSAP 논리 격리 및 규정 준수 재검증")
]

for row_idx, (cat, desc, soln) in enumerate(issues, start=1):
    table.cell(row_idx, 0).text = cat
    table.cell(row_idx, 1).text = desc
    table.cell(row_idx, 2).text = soln
    for col_idx in range(cols):
        cell = table.cell(row_idx, col_idx)
        if row_idx % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_BG_LIGHT
        set_font(cell.text_frame.paragraphs[0], size=12)


# -------------------------------------------------------------
# Slide 4: 공공 클래식 Fade-out 추진 계획 및 일정
# -------------------------------------------------------------
slide = prs.slides.add_slide(blank_layout)
add_header(slide, "3. 공공 클래식 우선 추진 로드맵", "정부 예산 편성 시기(8월) 이전 선제 공지 및 2027년 본격 이관")

# 타임라인을 시각화하는 4개 박스 배치
steps = [
    ("5월~6월", "전수 조사", "사용 고객, 상품 현황 파악\nClassic-VPC 기능 Gap 분석\n(이영대님, 조정수님 데이터 교차)"),
    ("7월", "내부 조율", "작업 계획 초안 수립\n내부 개발조직/Sales/TAM 사전 공유\n의사결정 사항 임원 보고"),
    ("8월 이전", "고객 공지 발송", "공공기관 예산 편성기 타겟\n내년도 예산 반영 유도\n고객 공지문 초안 작성"),
    ("2027년~", "본격 이관 시행", "과거 하이퍼바이저 이관 이력 활용\nNAS/CDB 마이그레이션 툴 투입\n본격적인 VPC 전환")
]

for i, (date, title, detail) in enumerate(steps):
    box_left = Inches(0.8 + i * 2.95)
    # 박스 배경
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, box_left, Inches(2.2), Inches(2.8), Inches(4.2))
    box.fill.solid()
    if i == 2: # 8월 공지 단계를 하이라이트
        box.fill.fore_color.rgb = RGBColor(235, 243, 255)
        box.line.color.rgb = COLOR_SECONDARY
        box.line.width = Pt(2)
    else:
        box.fill.fore_color.rgb = COLOR_BG_LIGHT
        box.line.color.rgb = RGBColor(210, 215, 225)
        
    tf = box.text_frame
    tf.word_wrap = True
    
    # 일정
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = date + "\n"
    set_font(r, size=16, bold=True, color=COLOR_SECONDARY)
    
    # 제목
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = title + "\n\n"
    set_font(r2, size=18, bold=True, color=COLOR_PRIMARY)
    
    # 상세내용
    p3 = tf.add_paragraph()
    r3 = p3.add_run()
    r3.text = detail
    set_font(r3, size=12, color=COLOR_TEXT)


# -------------------------------------------------------------
# Slide 5: 임원 의사결정 필요 사항 (수환님 보고 안건)
# -------------------------------------------------------------
slide = prs.slides.add_slide(blank_layout)
add_header(slide, "4. 핵심 경영진 의사결정 요청 안건", "인프라 투자 비용 및 고객 보전 정책에 대한 가이드라인 확립 필요")

agendas = [
    ("안건 1. 미국·독일 리전 VPC 신규 구축 여부", "• 현황: 글로벌 리전 VPC 부재로 이관 전면 중단\n• 리스크: 대형 고객(벨로프 등) 이탈 가능성\n• 결정 필요사항: 투자 비용 검토 후 승인 또는 글로벌 리전 예외 처리 여부"),
    ("안건 2. 이관에 따른 고객 비용 증가분 보전 책 책정", "• 현황: 상품 스펙 차이로 인해 VPC 이관 시 이용 요금 대폭 상승 케이스 발생\n• 결정 필요사항: 인상 금액을 회사가 일정 기간 보전할 것인가, 고객에게 청구할 것인가"),
    ("안건 3. 마이그레이션 수행 주체 및 MSP 비용 분담", "• 현황: OS/DB/어플리케이션 전반의 기술 지원 인력 부족\n• 결정 필요사항: 전문 MSP 솔루션/인력 활용 시 비용 부담 주체 지정 (회사 vs 고객)"),
    ("안건 4. 다운타임 발생에 대한 Credit 지급 기준", "• 현황: 대형 고객 이관 작업 시 서비스 중단(다운타임) 수반 불가피\n• 결정 필요사항: SLA 위반에 준하는 크레딧 보상 범위 및 예외 조항 마련")
]

for i, (title, detail) in enumerate(agendas):
    row = i // 2
    col = i % 2
    b_left = Inches(0.8 + col * 5.9)
    b_top = Inches(1.8 + row * 2.6)
    
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, b_left, b_top, Inches(5.6), Inches(2.3))
    box.fill.solid()
    box.fill.fore_color.rgb = COLOR_BG_LIGHT
    box.line.color.rgb = COLOR_SECONDARY
    
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.15)
    tf.margin_left = Inches(0.2)
    
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title + "\n"
    set_font(r, size=14, bold=True, color=COLOR_PRIMARY)
    
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = detail
    set_font(r2, size=11, color=COLOR_TEXT)


# -------------------------------------------------------------
# Slide 6: 전사 TF 구성 및 조직별 R&R
# -------------------------------------------------------------
slide = prs.slides.add_slide(blank_layout)
add_header(slide, "5. 전사 협업 Task Force 구성안", "성공적인 Fade-out을 위한 각 조직별 명확한 역할 정의")

# R&R 테이블 (행 7, 열 3)
rows, cols = 7, 3
left, top, width, height = Inches(0.8), Inches(1.8), Inches(11.733), Inches(4.8)
table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
table = table_shape.table

table.columns[0].width = Inches(2.2)
table.columns[1].width = Inches(3.5)
table.columns[2].width = Inches(6.033)

headers_rr = ["참여 조직", "주요 역할 (R&R)", "당면 액션 아이템"]
for i, h in enumerate(headers_rr):
    cell = table.cell(0, i)
    cell.text = h
    cell.fill.solid()
    cell.fill.fore_color.rgb = COLOR_PRIMARY
    set_font(cell.text_frame.paragraphs[0], size=14, bold=True, color=RGBColor(255,255,255))

rr_data = [
    ("PM 조직 (TVT)", "프로젝트 전체 총괄 및 매니징", "전수조사 내용 바탕 이슈 점검 및 7월 보고서 수립"),
    ("인프라 구축 조직", "VPC 이관 예정 물량 사전 선구축", "글로벌 리전(미·독) VPC 인프라 투자비용 가산출"),
    ("기획 및 개발 조직", "IaaS / PaaS / SaaS 상품 및 기능 개발", "Classic-VPC 기능 갭 보완 및 이관 자동화 툴 개발"),
    ("고객 접점 (Sales/TAM)", "VIP 및 대형 고객사 1:1 직접 대응", "고객사 수용 가능한 일정 협의 및 이탈 방지 영업"),
    ("보안 조직", "CSAP 인증 유지 및 논리 격리 검증", "이관 후 VPC 환경 보안 규정 및 가이드라인 제시"),
    ("CS 조직", "대고객 공지문 발송 및 인바운드 대응", "회원종류별(개인/기업/파트너) 문의 대응 가이드 준비")
]

for row_idx, (org, role, action) in enumerate(rr_data, start=1):
    table.cell(row_idx, 0).text = org
    table.cell(row_idx, 1).text = role
    table.cell(row_idx, 2).text = action
    for col_idx in range(cols):
        cell = table.cell(row_idx, col_idx)
        if row_idx % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_BG_LIGHT
        set_font(cell.text_frame.paragraphs[0], size=12)

# 파일 저장
prs.save("Classic_Fadeout_Strategy.pptx")
print("성공적으로 'Classic_Fadeout_Strategy.pptx' 파일이 생성되었습니다.")